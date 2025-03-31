import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os

# Ensure directories exist
os.makedirs("slides", exist_ok=True)
os.makedirs("output", exist_ok=True)

# Replace NEXT_PRESENTER text
def insert_next_presenter(slide, next_name):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if 'NEXT_PRESENTER' in run.text:
                        run.text = run.text.replace('NEXT_PRESENTER', next_name)

# Replace image placeholders with uploaded images
def replace_image_placeholders(slide, photo_file, logo_file):
    for shape in slide.shapes:
        if shape.name == "PhotoPlaceholder" and photo_file:
            slide.shapes.add_picture(photo_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "LogoPlaceholder" and logo_file:
            with Image.open(logo_file) as img:
                img_width, img_height = img.size
            placeholder_ratio = shape.width / shape.height
            img_ratio = img_width / img_height
            if img_ratio > placeholder_ratio:
                new_width = shape.width
                new_height = shape.width / img_ratio
            else:
                new_height = shape.height
                new_width = shape.height * img_ratio
            left = shape.left + (shape.width - new_width) / 2
            top = shape.top + (shape.height - new_height) / 2
            slide.shapes.add_picture(logo_file, left, top, new_width, new_height)

# Create individual slide with photo/logo/tag and BNIP ID
def create_individual_slide(name, company, field, photo_file, logo_file, green, gold):
    csv_file = 'bnipersons.csv'
    if os.path.exists(csv_file):
        df = pd.read_csv(csv_file)
    else:
        df = pd.DataFrame(columns=['ID', 'Name'])

    new_id = f"BNIP-{(len(df)+1):04d}"
    df = pd.concat([df, pd.DataFrame([[new_id, name]], columns=['ID', 'Name'])], ignore_index=True)
    df.to_csv(csv_file, index=False)

    prs = Presentation("slides/individual_template.pptx")
    slide = prs.slides[0]

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = run.text.replace("NAME", name)
                    run.text = run.text.replace("COMPANY_NME", company)
                    run.text = run.text.replace("FIELD", field)

    replace_image_placeholders(slide, photo_file, logo_file)

    if green:
        slide.shapes.add_picture("green_member_tag.png", Inches(0.9), Inches(4), width=Inches(3), height=Inches(0.5))
    if gold:
        slide.shapes.add_picture("gold_club_member_tag.png", Inches(0.9), Inches(4.6), width=Inches(3), height=Inches(0.7))

    output_path = f"output/{new_id}.pptx"
    prs.save(output_path)
    return new_id, output_path

# Create final presentation with theme slide, individual slides, and constant slide
def create_final_ppt(theme, names):
    df = pd.read_csv("bnipersons.csv")
    ids = df.set_index("Name").reindex(names)["ID"].tolist()

    final = Presentation()

    # Add theme slide(s)
    theme_prs = Presentation("slides/theme_slide.pptx")
    for slide in theme_prs.slides:
        final.slides.add_slide(slide)

    # Add individual slides
    for i, name in enumerate(names):
        pid = ids[i]
        next_name = names[(i + 1) % len(names)]
        slide_path = f"slides/{pid}.pptx"
        if not os.path.exists(slide_path):
            continue
        prs = Presentation(slide_path)
        for slide in prs.slides:
            insert_next_presenter(slide, next_name)
            final.slides.add_slide(slide)

    # Add constant slide(s)
    const_prs = Presentation("slides/constant_slide.pptx")
    for slide in const_prs.slides:
        final.slides.add_slide(slide)

    out_path = "output/final_presentation.pptx"
    final.save(out_path)
    return out_path

# ---------- Streamlit UI ----------
st.set_page_config(page_title="BNI Slide Generator", layout="centered")
st.title("📽️ BNI Slide Generator")

tab1, tab2 = st.tabs(["📊 Create Full Presentation", "👤 Create Individual Slide"])

with tab1:
    st.subheader("Create Weekly Presentation")
    theme = st.text_input("Enter theme of the week")
    excel_file = st.file_uploader("Upload Excel file (must include a 'Name' column)", type=["xlsx"])

    if st.button("Generate Final PPT") and excel_file:
        df = pd.read_excel(excel_file)
        if "Name" not in df.columns:
            st.error("❌ Excel must have a 'Name' column.")
        else:
            names = df["Name"].dropna().tolist()
            out_path = create_final_ppt(theme, names)
            st.success("✅ Final presentation created!")
            with open(out_path, "rb") as f:
                st.download_button("⬇️ Download Final PPT", f, file_name="final_presentation.pptx")

with tab2:
    st.subheader("Create Individual Slide")
    name = st.text_input("Name")
    company = st.text_input("Company")
    field = st.text_input("Field of Work")
    photo = st.file_uploader("Upload Photo", type=["jpg", "jpeg", "png"])
    logo = st.file_uploader("Upload Logo", type=["jpg", "jpeg", "png"])
    green = st.checkbox("Green Member")
    gold = st.checkbox("Gold Club Member")

    if st.button("Create Slide") and name and company and field and photo:
        sid, output = create_individual_slide(name, company, field, photo, logo, green, gold)
        st.success(f"✅ Slide created for {name} with ID {sid}")
        with open(output, "rb") as f:
            st.download_button("⬇️ Download Slide", f, file_name=f"{sid}.pptx")
